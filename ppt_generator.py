from pptx import Presentation
from pptx.util import Pt, Inches
import os
from typing import List, Dict, Any

def create_ppt(slides_data: List[Dict[str, Any]], output_ppt_path: str):
    """
    Creates a PowerPoint presentation from processed slide data.
    
    Args:
        slides_data: List of dicts, each containing:
            - 'clean_image_path': Path to the background image.
            - 'text_blocks': List of extracted text data.
            - 'original_dims': (width, height) tuple of original PDF page.
        output_ppt_path: Path to save the .pptx file.
    """
    prs = Presentation()
    
    # Set slide size based on the first slide's image aspect ratio if possible,
    # or just default to 16:9. For now, let's stick to standard 16:9 or 4:3.
    # We can adjust page size if needed, but let's keep it simple first.
    # prs.slide_width = Inches(13.333)
    # prs.slide_height = Inches(7.5)
    
    for slide_info in slides_data:
        blank_slide_layout = prs.slide_layouts[6] # 6 is usually blank
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 1. Add Background Image
        img_path = slide_info['clean_image_path']
        # We add the picture to cover the whole slide
        left = top = 0
        slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
        
        # 2. Add Text Boxes
        text_blocks = slide_info['text_blocks']
        
        # Dimensions for coordinate conversion
        start_w = prs.slide_width
        start_h = prs.slide_height
        
        for block in text_blocks:
            text = block.get('text', '')
            box = block.get('box_2d', [])
            
            if not text or not box:
                continue
                
            ymin, xmin, ymax, xmax = box
            
            # Convert 0-1000 to PPT units
            x = int((xmin / 1000) * start_w)
            y = int((ymin / 1000) * start_h)
            w = int((xmax - xmin) / 1000 * start_w)
            h = int((ymax - ymin) / 1000 * start_h)
            
            # Sanity check for small boxes
            if w <= 0 or h <= 0:
                continue
                
            txBox = slide.shapes.add_textbox(x, y, w, h)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(12) # Default size, difficult to guess exact from OCR
            # We could try to estimate font size based on box height if needed
            
    prs.save(output_ppt_path)
    print(f"Presentation saved: {output_ppt_path}")
